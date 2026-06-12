#!/usr/bin/env python3
"""
Build the GP2 defense deck (.pptx) — a god-mode, sales-narrative deck with hero stats,
the Confirmed-Tier funnel, real dashboard screenshots, a competitive table, and a
market/"the ask" arc.

Convert to .odp afterwards:
    libreoffice --headless --convert-to odp --outdir docs/ docs/ACR-QA_Defense.pptx
or simply: make deck

Usage:
    .venv/bin/python3 scripts/build_defense_deck.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ── Palette (KSIU + product) ────────────────────────────────────────────────
NAVY = RGBColor(0x00, 0x26, 0x54)
NAVY2 = RGBColor(0x00, 0x1A, 0x3D)
GOLD = RGBColor(0xC9, 0xA8, 0x4C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x14, 0x1A, 0x22)
GRAY = RGBColor(0x5B, 0x66, 0x72)
LIGHT = RGBColor(0xF4, 0xF6, 0xF9)
GREEN = RGBColor(0x18, 0x9E, 0x53)
RED = RGBColor(0xDC, 0x36, 0x36)
SLATE = RGBColor(0x27, 0x33, 0x42)

ROOT = Path(__file__).resolve().parent.parent
SHOTS = ROOT / "docs" / "presentation_assets" / "shots"
FIGS = ROOT / "ACR-QA-Book" / "figures"
LOGO = ROOT / "docs" / "presentation_assets" / "ksiu-logo.png"
OUT = ROOT / "docs" / "ACR-QA_Defense.pptx"

EMU_W, EMU_H = Inches(13.333), Inches(7.5)


# ── primitives ──────────────────────────────────────────────────────────────
def _bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _box(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    return tb.text_frame


def _run(p, text, size, color, *, bold=False, italic=False, font="Calibri"):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font
    return r


def _rect(slide, l, t, w, h, fill, *, line=None, rounded=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE, l, t, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _arrow(slide, l, t, w, h, color=GOLD):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, l, t, w, h)
    a.fill.solid()
    a.fill.fore_color.rgb = color
    a.line.fill.background()
    a.shadow.inherit = False
    return a


def _title(slide, text, *, sub=None, color=NAVY):
    tf = _box(slide, Inches(0.6), Inches(0.42), Inches(12.1), Inches(1.0))
    _run(tf.paragraphs[0], text, 30, color, bold=True)
    _rect(slide, Inches(0.62), Inches(1.28), Inches(1.7), Pt(4), GOLD)
    if sub:
        stf = _box(slide, Inches(0.62), Inches(1.34), Inches(12.0), Inches(0.5))
        _run(stf.paragraphs[0], sub, 14, GOLD, italic=True)
    return tf


def _center(p):
    p.alignment = PP_ALIGN.CENTER


# ── slide builders ──────────────────────────────────────────────────────────
def title_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, NAVY)
    _rect(s, 0, Inches(7.1), EMU_W, Inches(0.4), GOLD)
    if LOGO.exists():
        s.shapes.add_picture(str(LOGO), Inches(6.07), Inches(0.7), height=Inches(1.4))
    tf = _box(s, Inches(0.8), Inches(2.5), Inches(11.7), Inches(2.2))
    _center(tf.paragraphs[0])
    _run(tf.paragraphs[0], "ACR-QA", 62, WHITE, bold=True)
    p = tf.add_paragraph()
    _center(p)
    _run(p, "Automated Code Review & Quality Assurance", 24, GOLD, bold=True)
    p = tf.add_paragraph()
    _center(p)
    _run(
        p, "The trust layer for AI-written code — exploit-verified, cryptographically attested", 14, WHITE, italic=True
    )
    tf2 = _box(s, Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.4))
    for txt, sz, bold in [
        ("Ahmed Mahmoud Abbas", 18, True),
        ("Supervisor: Dr. Samy AbdelNabi", 13, False),
        ("Faculty of Computer Science & Engineering · CSE494 Graduation Project 2 · 2026", 11, False),
    ]:
        p = tf2.add_paragraph()
        _center(p)
        _run(p, txt, sz, WHITE, bold=bold)


def hook_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, NAVY2)
    tf = _box(s, Inches(1.0), Inches(1.5), Inches(11.3), Inches(4.5))
    _run(tf.paragraphs[0], "AI now writes a third of new code.", 40, WHITE, bold=True)
    p = tf.add_paragraph()
    p.space_before = Pt(16)
    _run(p, "And ", 26, RGBColor(0xC8, 0xD2, 0xE0))
    _run(p, "45% of AI-written code ships a known flaw", 26, GOLD, bold=True)
    _run(p, " (Veracode '25) —", 26, RGBColor(0xC8, 0xD2, 0xE0))
    p = tf.add_paragraph()
    _run(p, "while vulnerabilities per codebase jumped ", 26, RGBColor(0xC8, 0xD2, 0xE0))
    _run(p, "+107% in a year", 26, GOLD, bold=True)
    _run(p, " (Black Duck OSSRA '26).", 26, RGBColor(0xC8, 0xD2, 0xE0))
    p = tf.add_paragraph()
    p.space_before = Pt(28)
    _run(p, "Your scanner flags 1,900 issues. Which one breaches you?", 30, WHITE, bold=True)
    p = tf.add_paragraph()
    p.space_before = Pt(18)
    _run(
        p,
        "Nobody can review that. So teams ship blind — or pay $50k/year and still don't trust the output.",
        17,
        RGBColor(0x9F, 0xB0, 0xC4),
        italic=True,
    )


def problem_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    _title(s, "The Problem", sub="Three reasons teams can't trust automated code review today")
    cards = [
        (
            "Quality variance",
            "Manual review is inconsistent — expertise and attention vary. The SQL injection on line 47 slips through on the 50th PR of the day.",
            NAVY,
        ),
        (
            "Cost barrier",
            "Enterprise tools cost $10k–50k/year. Unaffordable for universities and startups — so they run a basic linter, or nothing.",
            GOLD,
        ),
        (
            "Hallucination",
            "AI explainers invent plausible-but-wrong guidance. A developer gets burned once and never trusts the tool again.",
            RED,
        ),
    ]
    x = Inches(0.6)
    for head, body, color in cards:
        _rect(s, x, Inches(2.0), Inches(3.95), Inches(4.2), LIGHT, rounded=True)
        _rect(s, x, Inches(2.0), Inches(3.95), Inches(0.16), color, rounded=False)
        tf = _box(s, x + Inches(0.25), Inches(2.35), Inches(3.45), Inches(3.6))
        _run(tf.paragraphs[0], head, 20, color, bold=True)
        p = tf.add_paragraph()
        p.space_before = Pt(10)
        _run(p, body, 14, INK)
        x = Emu(x + Inches(4.15))


def stat_cards_slide(prs, title, sub, cards, *, dark=False):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, NAVY if dark else WHITE)
    tcolor = WHITE if dark else NAVY
    tf = _box(s, Inches(0.6), Inches(0.5), Inches(12.1), Inches(1.0))
    _run(tf.paragraphs[0], title, 30, tcolor, bold=True)
    _rect(s, Inches(0.62), Inches(1.34), Inches(1.7), Pt(4), GOLD)
    if sub:
        stf = _box(s, Inches(0.62), Inches(1.4), Inches(12.0), Inches(0.5))
        _run(stf.paragraphs[0], sub, 14, GOLD, italic=True)
    n = len(cards)
    gap = Inches(0.3)
    total = Inches(12.1)
    cw = Emu(int((total - gap * (n - 1)) / n))
    x = Inches(0.6)
    for big, label in cards:
        _rect(s, x, Inches(2.4), cw, Inches(3.0), LIGHT if not dark else NAVY2, rounded=True)
        tf = _box(s, x, Inches(2.95), cw, Inches(2.0))
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        _center(p)
        _run(p, big, 46, GREEN if not dark else GOLD, bold=True)
        p2 = tf.add_paragraph()
        _center(p2)
        _run(p2, label, 13, INK if not dark else WHITE)
        x = Emu(x + cw + gap)


def bullets_slide(prs, title, bullets, *, sub=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    _title(s, title, sub=sub)
    body = _box(s, Inches(0.7), Inches(1.85), Inches(12.0), Inches(5.2))
    first = True
    for text, level, bold, color in bullets:
        p = body.paragraphs[0] if first else body.add_paragraph()
        first = False
        p.level = level
        _run(p, ("●  " if level == 0 else "–  ") + text, 19 - level * 3, color or INK, bold=bold)
        p.space_after = Pt(7)


def image_slide(prs, title, img, *, caption=None, sub=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    _title(s, title, sub=sub)
    if img.exists():
        # fit within content area 12.1 x 4.9 preserving aspect (img is 2880x1800 = 1.6)
        pic = s.shapes.add_picture(str(img), Inches(1.0), Inches(1.95), width=Inches(11.3))
        # frame
        fr = _rect(s, Inches(0.97), Inches(1.92), Inches(11.36), Emu(pic.height + Inches(0.06)), NAVY, rounded=False)
        fr.fill.background()
        fr.line.color.rgb = NAVY
        fr.line.width = Pt(1.5)
        spTree = s.shapes._spTree
        spTree.remove(pic._element)
        spTree.append(pic._element)
    if caption:
        ctf = _box(s, Inches(0.6), Inches(6.95), Inches(12.1), Inches(0.45))
        _run(ctf.paragraphs[0], caption, 12, GRAY, italic=True)


def full_image_slide(prs, title, img, *, sub=None, caption=None, dark=False):
    """Figure slide: title bar at top, image fills the remaining space edge-to-edge."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_color = NAVY if dark else WHITE
    _bg(s, bg_color)
    title_color = WHITE if dark else NAVY
    # compact title strip
    tf = _box(s, Inches(0.55), Inches(0.22), Inches(12.2), Inches(0.75))
    _run(tf.paragraphs[0], title, 24, title_color, bold=True)
    _rect(s, Inches(0.57), Inches(0.93), Inches(1.5), Pt(3), GOLD)
    if sub:
        stf = _box(s, Inches(0.57), Inches(0.97), Inches(12.0), Inches(0.38))
        _run(stf.paragraphs[0], sub, 12, GOLD, italic=True)
    # image: tall content zone 1.4" → 7.1" = 5.7" tall, full width
    img_top = Inches(1.42)
    img_h_max = Inches(5.7) if not caption else Inches(5.3)
    if img.exists():
        # add at full width first to get true aspect
        pic = s.shapes.add_picture(str(img), Inches(0.5), img_top, width=Inches(12.33))
        if pic.height > img_h_max:
            # too tall — constrain by height instead
            s.shapes._spTree.remove(pic._element)
            pic = s.shapes.add_picture(str(img), Inches(0.5), img_top, height=img_h_max)
            # centre horizontally
            pic.left = int((EMU_W - pic.width) / 2)
        # bring to front
        spTree = s.shapes._spTree
        spTree.remove(pic._element)
        spTree.append(pic._element)
    if caption:
        ctf = _box(s, Inches(0.55), Inches(6.88), Inches(12.2), Inches(0.52))
        _run(ctf.paragraphs[0], caption, 11, GRAY, italic=True)


def funnel_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    _title(s, "From Noise to Trust", sub="The Confirmed Tier: the few findings precise enough to auto-block a merge")
    stages = [
        ("1,942", "raw findings", "30–70% are noise (industry baseline)", SLATE, Inches(9.6)),
        ("219", "security-tier", "filtered by canonical rule set", NAVY, Inches(7.6)),
        ("151", "+ taint gate", "HTTP-source confirmed", RGBColor(0x1B, 0x4D, 0x89), Inches(5.8)),
        ("55", "Confirmed Tier", "96.4% precision · exploit-verified", GREEN, Inches(4.0)),
    ]
    y = Inches(2.05)
    for big, label, note, color, w in stages:
        left = Emu(int((EMU_W - w) / 2))
        _rect(s, left, y, w, Inches(0.92), color, rounded=True)
        tf = _box(s, left, y + Inches(0.04), w, Inches(0.85))
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        _center(p)
        _run(p, f"{big}   ", 26, WHITE, bold=True)
        _run(p, f"{label}", 15, WHITE)
        ntf = _box(s, Emu(left + w + Inches(0.15)), y + Inches(0.2), Inches(3.2), Inches(0.6))
        _run(ntf.paragraphs[0], note, 12, GRAY, italic=True)
        y = Emu(y + Inches(1.18))
    foot = _box(s, Inches(0.6), Inches(6.95), Inches(12.1), Inches(0.45))
    _run(
        foot.paragraphs[0],
        "Turn the 55 on as a required status check. The other ~1,887 stay visible and fixable — they just don't block the merge.",
        12,
        NAVY,
        italic=True,
    )


def eval_table_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    _title(
        s, "Evaluation & Results", sub="Not just built — measured on adversarial corpora, against the standard tools"
    )
    rows = [
        ("Confirmed-Tier precision (auto-block stratum)", "96.4%  (95% CI 90.9–100%)"),
        ("CVE recall (pre-registered battery)", "100% — 8 / 8 detectable"),
        ("Head-to-head F1", "98.2%  vs Semgrep 45.7% / Bandit 21.8%"),
        ("OWASP Top 10 coverage", "9 / 10 categories"),
        ("RealVuln 2026 leaderboard", "25.1% — beats Semgrep, Snyk, SonarQube"),
        ("Test suite / CORE coverage", "3,247 tests · 88%"),
    ]
    tbl = s.shapes.add_table(len(rows) + 1, 2, Inches(0.7), Inches(1.95), Inches(12.0), Inches(4.6)).table
    tbl.columns[0].width = Inches(6.8)
    tbl.columns[1].width = Inches(5.2)
    for c, h in enumerate(["What we measured", "Result"]):
        cell = tbl.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        r = cell.text_frame.paragraphs[0].runs[0]
        r.font.color.rgb = WHITE
        r.font.bold = True
        r.font.size = Pt(15)
    for i, (a, b) in enumerate(rows, 1):
        for c, val in enumerate((a, b)):
            cell = tbl.cell(i, c)
            cell.text = val
            r = cell.text_frame.paragraphs[0].runs[0]
            r.font.size = Pt(14)
            if c == 1:
                r.font.bold = True
                r.font.color.rgb = GREEN
            else:
                r.font.color.rgb = INK
    foot = _box(s, Inches(0.6), Inches(6.85), Inches(12.1), Inches(0.5))
    _run(
        foot.paragraphs[0],
        "Two corpora · pre-registered methodology · Wilson CIs · external Bandit gate (non-tautological).",
        12,
        GRAY,
        italic=True,
    )


def proofs_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    _title(s, "Why You Can Trust the Numbers")
    items = [
        (
            "We separated detection from trust",
            "Most tools emit 30–70% false positives. The Confirmed Tier hits 96.4% — precise enough to auto-block.",
            NAVY,
        ),
        (
            "We don't claim a vulnerability — we detonate it",
            "Real payloads in a Docker sandbox: SQLi ' OR 1=1, SSTI {{7*7}}→49, command injection. Verified live: exploits fire; safe code correctly does NOT.",
            GREEN,
        ),
        (
            "Every result is signed",
            "ECDSA-P256 + Sigstore Rekor + SLSA L3. An auditor verifies the exact scan and result in one command.",
            GOLD,
        ),
    ]
    y = Inches(1.9)
    for head, body, color in items:
        _rect(s, Inches(0.7), y, Inches(0.16), Inches(1.4), color, rounded=False)
        tf = _box(s, Inches(1.05), y, Inches(11.4), Inches(1.4))
        _run(tf.paragraphs[0], head, 21, color, bold=True)
        p = tf.add_paragraph()
        _run(p, body, 15, INK)
        y = Emu(y + Inches(1.62))


def competitive_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    _title(
        s,
        "Competitive Position",
        sub="The one quadrant the market leaves open: open-source, first-party, in-CI, attested, at $0",
    )
    headers = ["Capability", "Snyk", "Semgrep", "GHAS", "ACR-QA"]
    rows = [
        ("Exploit verification (Docker)", "✗", "✗", "✗", "✓"),
        ("Re-exploit to verify the fix", "✗", "✗", "✗", "✓"),
        ("Cryptographic attestation", "✗", "✗", "✗", "✓"),
        ("Confirmed Tier (auto-block)", "✗", "✗", "✗", "✓ 96.4%"),
        ("Self-hosted / $0 recurring", "✗", "✗", "✗", "✓"),
    ]
    tbl = s.shapes.add_table(len(rows) + 1, len(headers), Inches(0.7), Inches(2.0), Inches(12.0), Inches(4.0)).table
    tbl.columns[0].width = Inches(5.0)
    for c in range(1, 5):
        tbl.columns[c].width = Inches(1.75)
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        r = cell.text_frame.paragraphs[0].runs[0]
        r.font.color.rgb = WHITE if c < 4 else GOLD
        r.font.bold = True
        r.font.size = Pt(14)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER if c else PP_ALIGN.LEFT
    for i, row in enumerate(rows, 1):
        for c, val in enumerate(row):
            cell = tbl.cell(i, c)
            cell.text = val
            para = cell.text_frame.paragraphs[0]
            r = para.runs[0]
            r.font.size = Pt(14)
            if c == 0:
                r.font.color.rgb = INK
            else:
                para.alignment = PP_ALIGN.CENTER
                r.font.bold = True
                r.font.color.rgb = GREEN if val.startswith("✓") else RGBColor(0xC2, 0xC8, 0xCF)
            if c == 4 and val.startswith("✓"):
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xE9, 0xF7, 0xEF)


def section_slide(prs, title, subtitle):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, NAVY)
    _rect(s, 0, Inches(7.1), EMU_W, Inches(0.4), GOLD)
    tf = _box(s, Inches(0.8), Inches(2.8), Inches(11.7), Inches(2.0))
    _center(tf.paragraphs[0])
    _run(tf.paragraphs[0], title, 46, WHITE, bold=True)
    p = tf.add_paragraph()
    _center(p)
    _run(p, subtitle, 20, GOLD)


def the_ask_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, NAVY)
    tf = _box(s, Inches(0.9), Inches(1.0), Inches(11.5), Inches(1.4))
    _run(tf.paragraphs[0], "What ACR-QA Delivers", 32, WHITE, bold=True)
    _rect(s, Inches(0.92), Inches(1.85), Inches(1.7), Pt(4), GOLD)
    lines = [
        ("Trust", "96.4% Confirmed-Tier precision — high enough to auto-block a merge"),
        ("Proof", "Exploit-verified in a sandbox + cryptographically signed — not guesses"),
        ("Reach", "19 engines · Python / JS / Go · 9/10 OWASP · 100% CVE recall"),
        ("Price", "Self-hosted, your data never leaves, $0 recurring — vs $10–50k/year"),
    ]
    y = Inches(2.3)
    for k, v in lines:
        _rect(s, Inches(0.95), y, Inches(1.7), Inches(0.7), GOLD, rounded=True)
        ktf = _box(s, Inches(0.95), y + Inches(0.06), Inches(1.7), Inches(0.6))
        ktf.vertical_anchor = MSO_ANCHOR.MIDDLE
        _center(ktf.paragraphs[0])
        _run(ktf.paragraphs[0], k, 17, NAVY, bold=True)
        vtf = _box(s, Inches(2.9), y + Inches(0.04), Inches(9.6), Inches(0.7))
        vtf.vertical_anchor = MSO_ANCHOR.MIDDLE
        _run(vtf.paragraphs[0], v, 17, WHITE)
        y = Emu(y + Inches(0.92))
    foot = _box(s, Inches(0.95), Inches(6.3), Inches(11.5), Inches(0.8))
    _run(
        foot.paragraphs[0],
        "Detection you can trust, exploit-proof you can verify, a signed record of all of it — on any machine, at zero cost.",
        16,
        GOLD,
        bold=True,
        italic=True,
    )


def B(text, level=0, bold=False, color=None):
    return (text, level, bold, color)


def build():
    """
    18-slide photo-first deck.

    Slide map:
     1  Title
     2  Hook (dark, text)
     3  Problem — 3 cards
     4  Market reality — stat cards (dark)
     5  Solution — bullets (what ACR-QA does)
     6  FIGURE: PR_OPERATING_POINTS — "two modes, one scan" (kills numbers question)
     7  FIGURE: arch_overview — full system diagram (replaces architecture bullets)
     8  Live Dashboard overview screenshot
     9  FIGURE: FUNNEL_SLIDE — 1,942 → 55 (real figure, not generated)
    10  FIGURE: HEAD_TO_HEAD — bar chart vs Semgrep / Bandit
    11  FIGURE: REALVULN_LEADERBOARD — real-world beats everyone
    12  FIGURE: verified_remediation — 4-phase exploit chain (the "wow" slide)
    13  Section divider: Live Demo
    14  Dashboard: run-detail screenshot
    15  Dashboard: OWASP compliance screenshot
    16  Dashboard: attestation screenshot (Signature Verified)
    17  Competitive table
    18  The Ask (dark, closing sell)
    """
    prs = Presentation()
    prs.slide_width, prs.slide_height = EMU_W, EMU_H

    # 1 — Title
    title_slide(prs)

    # 2 — Hook
    hook_slide(prs)

    # 3 — Problem
    problem_slide(prs)

    # 4 — Market reality
    stat_cards_slide(
        prs,
        "The Market Reality",
        "Why this matters — now",
        [
            ("$10–50k", "per year for enterprise SAST tools"),
            ("+107%", "open-source vulns / codebase in 2024  (Black Duck OSSRA 2026)"),
            ("45%", "of AI-written code ships a known flaw  (Veracode 2025)"),
            ("$0", "ACR-QA — self-hosted, your data never leaves"),
        ],
        dark=True,
    )

    # 5 — Solution bullets
    bullets_slide(
        prs,
        "The Solution: ACR-QA",
        [
            B("A trust layer on top of your scanners — answers one question at merge time:", 0, True, NAVY),
            B('"Is this finding real enough to block automatically?"', 1, False, GREEN),
            B("RAG-grounded AI — cites the rule, entropy filter rejects hallucinated explanations"),
            B("Confirmed Tier — 4-gate filter at 96.4% precision, safe to auto-block without human review"),
            B("Exploit verification — fires a real payload in an isolated Docker sandbox"),
            B("ECDSA-P256 + post-quantum attestation — every scan signed, tamper-evident"),
            B("Self-hosted · zero recurring cost · full data privacy", 0, True, NAVY),
        ],
    )

    # 6 — PR operating points (precision-recall scatter — kills "which number?" question)
    full_image_slide(
        prs,
        "Two Operating Modes — One Scan",
        FIGS / "PR_OPERATING_POINTS.png",
        sub="Confirmed Tier (auto-block) vs Full Output (triage) — same pipeline, different strictness",
        caption=(
            "Top-right is the ideal zone. Confirmed Tier (gold diamond, 96.4% precision) is safe to auto-block. "
            "Full output (navy circle, 91% recall) is the developer triage view. "
            "Both are real — they're two points on the same precision-recall curve."
        ),
    )

    # 7 — Architecture: real diagram (replaces bullet list)
    full_image_slide(
        prs,
        "System Architecture",
        FIGS / "arch_overview.png",
        sub="19 engines · 3 language adapters · 12-stage async pipeline · 52-endpoint FastAPI",
        caption=(
            "Push triggers a Celery worker: 12 tools run in parallel, every output normalised to CanonicalFinding, "
            "trust gates applied, results signed and persisted. End-to-end: 14–90 s."
        ),
    )

    # 8 — Live dashboard overview
    image_slide(
        prs,
        "Live Dashboard — Fleet Overview",
        SHOTS / "overview.png",
        sub="Real scans on real repos — psf/requests · encode/httpx · FastAPI · Flask",
        caption=(
            "Confirmed Tier tile fetches live from the API — never estimated. "
            "Trust Layer banner shows precision, recall, F1, and self-scan result in real time."
        ),
    )

    # 9 — Precision funnel (real thesis figure)
    full_image_slide(
        prs,
        "The Precision Funnel — From Noise to Trust",
        FIGS / "FUNNEL_SLIDE.png",
        sub="24-repo adversarial corpus · 1,942 raw findings → 55 Confirmed Tier · CVE recall preserved at every rung",
        caption=(
            "Note: these numbers are the full evaluation corpus (24 repos). "
            "The live demo scans one app and produces proportionally fewer confirmed findings. "
            "Same filter — smaller input."
        ),
    )

    # 10 — Head-to-head bar chart (real thesis figure)
    full_image_slide(
        prs,
        "Head-to-Head Benchmark — Precision · CVE Recall · F₁",
        FIGS / "HEAD_TO_HEAD.png",
        sub="Same 30-repo adversarial corpus · same 8-CVE recall battery · conservative triage",
        caption=(
            "ACR-QA P4 (Confirmed Tier): 96% precision · 100% CVE recall · F₁ = 98.2% — "
            ">52 pp margin over Semgrep CE. Source: docs/evaluation/HEAD_TO_HEAD_BENCHMARK.md"
        ),
    )

    # 11 — RealVuln leaderboard (real thesis figure)
    full_image_slide(
        prs,
        "RealVuln 2026 — Real-World Recall Leaderboard",
        FIGS / "REALVULN_LEADERBOARD.png",
        sub="22 real-world Python CVE apps · third-party ground truth (arXiv:2604.13764) · strict CWE+file+line±10 matching",
        caption=(
            "ACR-QA Full Output: 25.1% beats Bandit (19.4%), Semgrep CE (17.5%), Snyk (17.4%), SonarQube (6.5%). "
            "+LLM augmented: 32.4%. ACR-QA detectable-subset recall: 37.8%."
        ),
    )

    # 12 — Verified remediation flow (the "wow" slide before demo)
    full_image_slide(
        prs,
        "Exploit Verification — Proven, Not Claimed",
        FIGS / "verified_remediation.png",
        sub="4-phase chain: detect → detonate → patch → re-detonate · ECDSA-signed bundle at every stage",
        caption=(
            "Phase 1: rule→exploit category mapping (13 categories). "
            "Phase 2: Docker sandbox fires real payload (SQLi ' OR 1=1, CMDi ; echo PWNED, SSTI {{7*7}}). "
            "Phase 3: AI patch applied, exploit re-fired → must fail. "
            "Phase 4: (vuln_proof, fix_diff, fix_proof) signed as one attestation bundle."
        ),
    )

    # 13 — Live demo divider
    section_slide(prs, "Live Demo", "Real repo · real finding · real exploit · signed receipt")

    # 14 — Run detail
    image_slide(
        prs,
        "Run Detail — payments-api",
        SHOTS / "run-detail.png",
        sub="64 findings · 13 HIGH · 4 Confirmed Tier — SQL injection, eval, hardcoded secrets",
        caption=(
            "Each finding: canonical rule ID, severity, confidence score, taint-confirmed flag. "
            "Tabs: compliance mapping, attestation, PR risk score."
        ),
    )

    # 15 — OWASP compliance
    image_slide(
        prs,
        "OWASP Top 10 — Compliance View",
        SHOTS / "compliance.png",
        sub="A03 Injection: 8 findings · A02 Crypto: 3 · 9/10 categories covered",
        caption=(
            "Score calculated from passed categories. On numpy, pandas, pydantic, requests — "
            "0 HIGH findings reported (0.0% false-positive rate on clean, mature code)."
        ),
    )

    # 16 — Attestation
    image_slide(
        prs,
        "Cryptographic Attestation — Signature Verified",
        SHOTS / "attestation.png",
        sub="ECDSA-P256 primary · Dilithium3 post-quantum · public key embedded in bundle",
        caption=(
            "Tamper-evident: alter any finding and verification fails instantly. "
            "Verify in one command: python scripts/verify_attestation.py <bundle.json>. "
            "EU CRA Sept 2026 requires this class of provenance."
        ),
    )

    # 17 — Competitive table
    competitive_slide(prs)

    # 18 — The Ask (closing sell)
    the_ask_slide(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    n = len(prs.slides._sldIdLst)
    print(f"✓ Deck written: {OUT}  ({n} slides)")


if __name__ == "__main__":
    build()
