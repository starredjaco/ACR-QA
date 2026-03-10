import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    # Create presentation
    prs = Presentation()
    
    # Define slide layouts
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]

    # Slide 1: Title
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "ACR-QA Phase 2 Progress Update"
    subtitle.text = "Competitive Features, Deep-Code Testing, and Production Readiness\nAhmed"

    # Slide 2: Context
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Quick Refresher: What is ACR-QA?"
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "A multi-tool, RAG-enhanced code quality platform"
    p = tf.add_paragraph()
    p.text = "Explains exactly WHY code is bad and HOW to fix it using AI"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Zero recurring cost for teams"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Currently integrating 7 distinct static analysis engines"
    p.level = 1

    # Slide 3: The Leap
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "The Leap (Prototype vs. Enterprise)"
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "Jan 18th Status (v2.0)"
    p = tf.add_paragraph()
    p.text = "Working prototype (Detection, RAG AI, basic Dashboard)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Today's Status (v2.7)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Transformed into an enterprise-grade CI/CD product"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Production Hardening & Reliability (273 tests)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "DevOps & Monitoring (GitHub Actions, Prometheus/Grafana)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Competitive 'God-Mode' Features"
    p.level = 1

    # Slide 4: Enterprise Architecture
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Enterprise Architecture & Policy-as-Code"
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "Quality Gates"
    p = tf.add_paragraph()
    p.text = "Teams can define pass/fail thresholds programmatically in .acrqa.yml"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Inline Suppression"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Developers can ignore false positives directly in code via comments"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Language Adapters"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Restructured the core to easily plug in JavaScript/TypeScript next"
    p.level = 1

    # Slide 5: Feature 1
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Feature: The Test Gap Analyzer"
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "The Problem"
    p = tf.add_paragraph()
    p.text = "Tools show 80% coverage, but don't flag which complex functions are missed"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Our Innovation (ACR-QA)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Built an AST (Abstract Syntax Tree) parser to scan Python code physically"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Identifies exact untested functions sorted by logical complexity"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Fails the CI/CD pipeline if critical business logic skips tests"
    p.level = 1

    # Slide 6: Feature 2
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Feature: Security & Policy Automation"
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "OWASP Top 10 Compliance"
    p = tf.add_paragraph()
    p.text = "Every security finding is automatically mapped to OWASP categories and CWE IDs"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Automated Compliance Reporting generation for pipelines"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Policy Engine"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Built a robust validator for .acrqa.yml to prevent misconfigurations"
    p.level = 1

    # Slide 7: Feature 3
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Feature: Developer Fatigue Tuner"
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "The Problem"
    p = tf.add_paragraph()
    p.text = "Developers abandon static analysis due to false positive spam"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Our Innovation (Feedback Tuner)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Engine learns from 'false positive' votes on the dashboard"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Automatically calculates false-positive rates per rule"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Dynamically dials down severity of noisy rules without manual config"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Filters low-confidence AI explanations (min score: 0.7)"
    p.level = 1

    # Slide 8: Reliability
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Military-Grade Reliability (v2.7)"
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "God-Mode Testing Architecture"
    p = tf.add_paragraph()
    p.text = "Scaled from ~30 basic tests to 273 rigorous backend tests"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Full suite executes logically in under 6 seconds"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Simulated extreme loads: corrupt DBs, missing files, 1000-finding streams"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Zero crashes. 100% CI/CD pass rate on GitHub Actions."
    p.level = 1

    # Slide 9: Professionalization
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "GitHub Professionalization"
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "Ready for the Real World"
    p = tf.add_paragraph()
    p.text = "Auto-posting inline code review comments on live Pull Requests"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Official v2.7.0 Release Page published with comprehensive notes"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Phase 2 Roadmap fully structured and labeled via GitHub Issues"
    p.level = 1

    # Slide 10: Next Steps
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Roadmap (Phase 2 Launch)"
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "Language Expansion"
    p = tf.add_paragraph()
    p.text = "Launching JavaScript/TypeScript parsing this month"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Scientific Validation"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Conducting user studies + Precision/Recall bench tests"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Cloud Readiness"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Finalizing cloud deployment (AWS/GCP) for live demonstration"
    p.level = 1

    prs.save("docs/ACR_QA_Phase2_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()
    print("Created docs/ACR_QA_Phase2_Presentation.pptx")
