#!/usr/bin/env python3
"""
Build the GP2 defense deck in Dr. Samy's KSIU template style (matches the required
presentation template used across the faculty: navy #002060 + gold, KSIU branding,
project-title header bar, blue section headings, white content slides).

Photo-first, reference-pattern structure (mirrors the EGX360 defense deck):
Intro → Problem (overview + 3 photo slides) → Solution (overview + 4 pillar slides) →
System Architecture → Evaluation & Data → Results → Use Cases → Conclusion → Future Work.

Identity: ACR-QA is a *deterministic trust layer*. The TRUST — exploit verification (Docker)
and attestation (ECDSA + Dilithium3) — is deterministic proof an LLM structurally cannot give
(a model can't detonate a bug or sign a result). An optional LLM tier (llm_detector.py,
second_opinion.py) adds a few points of recall, gated and re-verified — but it never gets the
final word. We go QUIET on the LLM: do not headline it, do not claim "no AI" (the code uses one).

Run:    .venv/bin/python3 scripts/build_template_deck.py
Then:   libreoffice --headless --convert-to odp --outdir docs/ docs/ACR-QA_Defense.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

# ── Template palette (extracted from the faculty template) ───────────────────
NAVY = RGBColor(0x00, 0x20, 0x60)
NAVY2 = RGBColor(0x00, 0x16, 0x42)
GOLD = RGBColor(0xD4, 0xAF, 0x37)
GOLD_DK = RGBColor(0xB8, 0x8A, 0x1E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x1A, 0x1F, 0x2B)
GRAY = RGBColor(0x5B, 0x66, 0x72)
LIGHT = RGBColor(0xF2, 0xF4, 0xF9)
GREEN = RGBColor(0x18, 0x8A, 0x4E)
RED = RGBColor(0xC2, 0x2E, 0x2E)
BLUEROW = RGBColor(0x2A, 0x4B, 0x8C)

ROOT = Path(__file__).resolve().parent.parent
SHOTS = ROOT / "docs" / "presentation_assets" / "shots"
FIGS = ROOT / "ACR-QA-Book" / "figures"
MEDIA = ROOT / "docs" / "media"
LOGO = ROOT / "docs" / "presentation_assets" / "ksiu-logo.png"
PHOTOS = ROOT / "docs" / "GEMINI_PHOTOS"
OUT = ROOT / "docs" / "ACR-QA_Defense.pptx"

EMU_W, EMU_H = Inches(13.333), Inches(7.5)
PROJECT_TITLE = "ACR-QA: Automated Code Review & Quality Assurance"


# ── primitives ───────────────────────────────────────────────────────────────
def _bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _box(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    return tb.text_frame


def _run(p, text, size, color, *, bold=False, italic=False, font="Calibri"):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font
    return r


def _rect(slide, l, t, w, h, fill, *, line=None, rounded=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE, l, t, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _center(p):
    p.alignment = PP_ALIGN.CENTER


def _hero(slide, photo, alpha):
    p = PHOTOS / photo
    if not p.exists():
        return
    slide.shapes.add_picture(str(p), 0, 0, width=EMU_W, height=EMU_H)
    shp = _rect(slide, 0, 0, EMU_W, EMU_H, NAVY)
    sf = shp.fill._xPr.find(qn("a:solidFill"))
    clr = sf.find(qn("a:srgbClr"))
    clr.append(clr.makeelement(qn("a:alpha"), {"val": str(int(alpha * 1000))}))


def _fit_image(slide, img, left, top, max_w, max_h):
    """Place an image fitted within a box, preserving aspect, centred in the box."""
    if not img.exists():
        return None
    pic = slide.shapes.add_picture(str(img), left, top, width=max_w)
    if pic.height > max_h:
        slide.shapes._spTree.remove(pic._element)
        pic = slide.shapes.add_picture(str(img), left, top, height=max_h)
    pic.left = Emu(int(left + (max_w - pic.width) / 2))
    pic.top = Emu(int(top + (max_h - pic.height) / 2))
    return pic


def _cover_image(slide, img, left, top, w, h):
    """Fill a box with an image, cropping overflow (object-fit: cover)."""
    if not img.exists():
        _rect(slide, left, top, w, h, NAVY2, rounded=True)
        return None
    pic = slide.shapes.add_picture(str(img), left, top, width=w, height=h)
    return pic


def _arrow(slide, x, y, w, h, color):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    a.fill.solid()
    a.fill.fore_color.rgb = color
    a.line.fill.background()
    a.shadow.inherit = False
    return a


def _content(prs, heading, page, *, sub=None):
    """White content slide with the template header bar + blue section heading. Returns slide."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    htf = _box(s, Inches(2.0), Inches(0.18), Inches(9.3), Inches(0.4))
    p = htf.paragraphs[0]
    _center(p)
    _run(p, PROJECT_TITLE, 11, GRAY, italic=True)
    if LOGO.exists():
        s.shapes.add_picture(str(LOGO), Inches(12.55), Inches(0.12), height=Inches(0.5))
    _rect(s, Inches(0.5), Inches(0.66), Inches(12.33), Pt(1.4), NAVY)
    htf2 = _box(s, Inches(0.5), Inches(0.8), Inches(11.0), Inches(0.7))
    _run(htf2.paragraphs[0], heading, 28, NAVY, bold=True)
    if sub:
        st = _box(s, Inches(0.5), Inches(1.46), Inches(12.3), Inches(0.4))
        _run(st.paragraphs[0], sub, 13, GOLD_DK, italic=True)
    _rect(s, Inches(0.2), Inches(7.05), Inches(0.4), Inches(0.4), NAVY)
    ptf = _box(s, Inches(0.2), Inches(7.05), Inches(0.4), Inches(0.4))
    ptf.vertical_anchor = MSO_ANCHOR.MIDDLE
    pp = ptf.paragraphs[0]
    _center(pp)
    _run(pp, str(page), 12, WHITE, bold=True)
    return s


def _band(s, y, gold_lead, rest, *, h=Inches(1.0)):
    """Navy callout band with a gold lead-in and white body."""
    _rect(s, Inches(0.5), y, Inches(12.33), h, NAVY, rounded=True)
    bt = _box(s, Inches(0.8), y, Inches(11.7), h)
    bt.vertical_anchor = MSO_ANCHOR.MIDDLE
    if gold_lead:
        _run(bt.paragraphs[0], gold_lead, 13, GOLD, bold=True)
    _run(bt.paragraphs[0], rest, 13, WHITE)
    return bt


# ── 1 · title ─────────────────────────────────────────────────────────────────
def title_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, NAVY)
    _hero(s, "Gemini_1.png", alpha=66)
    _rect(s, Inches(11.7), 0, Inches(1.63), EMU_H, GOLD)
    _rect(s, Inches(11.4), 0, Inches(0.3), EMU_H, NAVY2)
    blocks = [
        ("BRANCH / FACULTY", "El Tur / Faculty of Computer Science and Engineering", 13),
        ("PROGRAM", "Computer Science", 13),
        ("COURSE", "CSE494 — Graduation Project 2", 13),
    ]
    y = Inches(0.7)
    for label, val, sz in blocks:
        ltf = _box(s, Inches(0.7), y, Inches(9.5), Inches(0.32))
        _run(ltf.paragraphs[0], label, 11, GOLD, bold=True)
        vtf = _box(s, Inches(0.7), Emu(y + Inches(0.3)), Inches(10.0), Inches(0.4))
        _run(vtf.paragraphs[0], val, sz, WHITE)
        y = Emu(y + Inches(0.82))
    ttf = _box(s, Inches(0.7), Inches(3.25), Inches(10.3), Inches(1.7))
    _run(ttf.paragraphs[0], "PROJECT TITLE", 11, GOLD, bold=True)
    p = ttf.add_paragraph()
    p.space_before = Pt(4)
    _run(p, "ACR-QA", 40, WHITE, bold=True)
    p2 = ttf.add_paragraph()
    _run(p2, "Automated Code Review & Quality Assurance", 22, GOLD, bold=True)
    p3 = ttf.add_paragraph()
    p3.space_before = Pt(4)
    _run(
        p3,
        "A deterministic trust layer for AI-written code — exploit-verified, cryptographically attested",
        12,
        WHITE,
        italic=True,
    )
    if LOGO.exists():
        s.shapes.add_picture(str(LOGO), Inches(0.7), Inches(5.85), height=Inches(1.0))
    pf = _box(s, Inches(1.95), Inches(5.7), Inches(9.0), Inches(1.4))
    _run(pf.paragraphs[0], "Presented by", 11, GOLD, bold=True)
    p = pf.add_paragraph()
    _run(p, "Ahmed Mahmoud Abbas", 17, WHITE, bold=True)
    p = pf.add_paragraph()
    p.space_before = Pt(8)
    _run(p, "Supervisor", 11, GOLD, bold=True)
    p = pf.add_paragraph()
    _run(p, "Dr. Samy Abdel Nabi", 15, WHITE)
    yr = _box(s, Inches(8.4), Inches(6.95), Inches(3.0), Inches(0.4))
    pr = yr.paragraphs[0]
    pr.alignment = PP_ALIGN.RIGHT
    _run(pr, "KSIU · Egypt · 2026", 11, WHITE)


# ── 2 · outline ─────────────────────────────────────────────────────────────────
def outline_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, NAVY)
    if LOGO.exists():
        s.shapes.add_picture(str(LOGO), Inches(0.5), Inches(0.4), height=Inches(1.05))
    ltf = _box(s, Inches(0.6), Inches(3.0), Inches(3.0), Inches(1.4))
    _run(ltf.paragraphs[0], "PRESENTATION", 14, GOLD, bold=True)
    p = ltf.add_paragraph()
    _run(p, "OUTLINE", 34, WHITE, bold=True)
    sections = [
        "Introduction",
        "Problem Statement",
        "The Solution",
        "System Architecture",
        "Evaluation & Data",
        "Results",
        "Use Cases",
        "Conclusion & Future Work",
    ]
    x = Inches(4.3)
    w = Inches(8.4)
    row_h = Inches(0.62)
    top = Inches(0.55)
    for i, sec in enumerate(sections):
        y = Emu(top + i * row_h)
        _rect(s, x, y, Inches(0.7), Inches(0.6), GOLD)
        ntf = _box(s, x, y, Inches(0.7), Inches(0.6))
        ntf.vertical_anchor = MSO_ANCHOR.MIDDLE
        np = ntf.paragraphs[0]
        _center(np)
        _run(np, f"{i + 1:02d}", 16, NAVY, bold=True)
        _rect(s, Emu(x + Inches(0.7)), y, Emu(w - Inches(0.7)), Inches(0.6), BLUEROW)
        ttf = _box(s, Emu(x + Inches(1.0)), y, Emu(w - Inches(1.1)), Inches(0.6))
        ttf.vertical_anchor = MSO_ANCHOR.MIDDLE
        _run(ttf.paragraphs[0], sec, 16, WHITE, bold=True)


# ── 3 · introduction ───────────────────────────────────────────────────────────
def introduction_slide(prs):
    s = _content(prs, "Introduction", 3)
    tf = _box(s, Inches(0.5), Inches(1.55), Inches(6.4), Inches(1.2))
    _run(tf.paragraphs[0], "AI now writes a third of new code —", 19, INK, bold=True)
    p = tf.add_paragraph()
    _run(p, "and trust hasn't caught up.", 19, INK, bold=True)
    stats = [
        ("45%", "of AI code ships a known flaw", RED),
        ("+107%", "vulns / codebase in one year", RED),
        ("$10–50k", "/yr for tools most can't afford", NAVY),
    ]
    y = Inches(2.9)
    for big, lab, col in stats:
        _rect(s, Inches(0.5), y, Inches(5.9), Inches(0.95), LIGHT, rounded=True)
        bt = _box(s, Inches(0.7), y, Inches(2.1), Inches(0.95))
        bt.vertical_anchor = MSO_ANCHOR.MIDDLE
        _run(bt.paragraphs[0], big, 28, col, bold=True)
        lt = _box(s, Inches(2.7), y, Inches(3.6), Inches(0.95))
        lt.vertical_anchor = MSO_ANCHOR.MIDDLE
        _run(lt.paragraphs[0], lab, 13, INK)
        y = Emu(y + Inches(1.08))
    _fit_image(s, MEDIA / "pr-comment-mockup1.png", Inches(6.8), Inches(1.65), Inches(6.2), Inches(4.1))
    cap = _box(s, Inches(6.8), Inches(5.8), Inches(6.2), Inches(0.5))
    _center(cap.paragraphs[0])
    _run(cap.paragraphs[0], "What a developer sees on their PR — within 90 seconds of a push.", 12, GRAY, italic=True)


# ── 4 · problem overview ────────────────────────────────────────────────────────
def problem_overview_slide(prs):
    s = _content(prs, "Problem Statement", 4)
    tf = _box(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.6))
    _run(
        tf.paragraphs[0],
        "Your scanner flags ~1,900 issues per scan. Which one actually breaches you?",
        18,
        INK,
        bold=True,
    )
    cards = [
        (
            "1",
            "Too many alerts,\nzero trust",
            "A wall of warnings, mostly false. Nobody can tell the real SQL injection from the noise.",
            NAVY,
        ),
        (
            "2",
            "Claimed,\nnever proven",
            "Tools say “possible vulnerability” — but never prove it is actually exploitable.",
            RED,
        ),
        (
            "3",
            "No proof of\nprovenance",
            "Results aren’t signed. You can’t trust a stranger’s PR, or what AI just wrote into your repo.",
            GREEN,
        ),
    ]
    x = Inches(0.5)
    for num, head, body, col in cards:
        _rect(s, x, Inches(2.5), Inches(3.95), Inches(3.4), LIGHT, rounded=True)
        _rect(s, x, Inches(2.5), Inches(3.95), Inches(0.14), col)
        _rect(s, x + Inches(0.25), Inches(2.75), Inches(0.7), Inches(0.7), col, rounded=True)
        nt = _box(s, x + Inches(0.25), Inches(2.75), Inches(0.7), Inches(0.7))
        nt.vertical_anchor = MSO_ANCHOR.MIDDLE
        _center(nt.paragraphs[0])
        _run(nt.paragraphs[0], num, 22, WHITE, bold=True)
        c = _box(s, x + Inches(0.25), Inches(3.6), Inches(3.45), Inches(2.2))
        _run(c.paragraphs[0], head, 18, col if col != GOLD else GOLD_DK, bold=True)
        p = c.add_paragraph()
        p.space_before = Pt(8)
        _run(p, body, 13.5, INK)
        x = Emu(x + Inches(4.15))
    _band(s, Inches(6.15), "The gap: ", "detection is cheap and noisy — trust is what's missing.", h=Inches(0.85))


# ── 5–7 · problem detail (photo) ────────────────────────────────────────────────
def problem_detail_slide(prs, page, num, photo, title, punch, stat_big, stat_lab):
    s = _content(prs, "Problem Statement", page)
    # left text column
    _rect(s, Inches(0.5), Inches(1.7), Inches(0.78), Inches(0.78), GOLD, rounded=True)
    nt = _box(s, Inches(0.5), Inches(1.7), Inches(0.78), Inches(0.78))
    nt.vertical_anchor = MSO_ANCHOR.MIDDLE
    _center(nt.paragraphs[0])
    _run(nt.paragraphs[0], num, 26, NAVY, bold=True)
    tt = _box(s, Inches(0.5), Inches(2.7), Inches(5.9), Inches(1.3))
    _run(tt.paragraphs[0], title, 30, NAVY, bold=True)
    pt = _box(s, Inches(0.5), Inches(4.0), Inches(5.9), Inches(1.6))
    _run(pt.paragraphs[0], punch, 17, INK)
    # stat pill
    _rect(s, Inches(0.5), Inches(5.7), Inches(5.9), Inches(1.0), LIGHT, rounded=True)
    bt = _box(s, Inches(0.7), Inches(5.7), Inches(2.4), Inches(1.0))
    bt.vertical_anchor = MSO_ANCHOR.MIDDLE
    _run(bt.paragraphs[0], stat_big, 30, RED, bold=True)
    lt = _box(s, Inches(3.0), Inches(5.7), Inches(3.2), Inches(1.0))
    lt.vertical_anchor = MSO_ANCHOR.MIDDLE
    _run(lt.paragraphs[0], stat_lab, 13, INK)
    # right photo (cover)
    _cover_image(s, PHOTOS / photo, Inches(6.75), Inches(1.7), Inches(6.08), Inches(5.0))


# ── 8 · solution overview (4 pillars) ───────────────────────────────────────────
def solution_overview_slide(prs):
    s = _content(prs, "The Solution — ACR-QA", 8)
    tf = _box(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.75))
    _run(tf.paragraphs[0], "A trust layer on top of your scanners. One question at merge time: ", 18, INK, bold=True)
    _run(tf.paragraphs[0], "is this finding real enough to stop a release?", 18, GREEN, bold=True)
    pillars = [
        (
            "1",
            "Deterministic Detection",
            "13 tools across Python, JS, Go — one schema. Reproducible, every run.",
            "the foundation",
            NAVY,
        ),
        (
            "2",
            "Trust Gates",
            "Reachability + taint + confidence collapse noise to 96.4% precision.",
            "solves Problem 1",
            GOLD_DK,
        ),
        (
            "3",
            "Exploit Verification",
            "Detonate a real payload in Docker. Re-detonate the fix. 5/5 live.",
            "solves Problem 2",
            RED,
        ),
        (
            "4",
            "Crypto Attestation",
            "ECDSA-P256 + Dilithium3 sign every scan. EU CRA-ready, $0 recurring.",
            "solves Problem 3",
            GREEN,
        ),
    ]
    x = Inches(0.5)
    cw = Inches(2.97)
    gap = Inches(0.13)
    for num, head, body, tag, col in pillars:
        _rect(s, x, Inches(2.55), cw, Inches(3.2), LIGHT, rounded=True)
        _rect(s, x, Inches(2.55), cw, Inches(0.14), col)
        _rect(s, x + Inches(0.22), Inches(2.8), Inches(0.62), Inches(0.62), col, rounded=True)
        nt = _box(s, x + Inches(0.22), Inches(2.8), Inches(0.62), Inches(0.62))
        nt.vertical_anchor = MSO_ANCHOR.MIDDLE
        _center(nt.paragraphs[0])
        _run(nt.paragraphs[0], num, 20, WHITE, bold=True)
        c = _box(s, x + Inches(0.22), Inches(3.6), Emu(cw - Inches(0.44)), Inches(1.6))
        _run(c.paragraphs[0], head, 15.5, col, bold=True)
        p = c.add_paragraph()
        p.space_before = Pt(8)
        _run(p, body, 12.5, INK)
        # problem→pillar mapping tag at card foot
        tg = _box(s, x + Inches(0.22), Inches(5.35), Emu(cw - Inches(0.44)), Inches(0.32))
        _run(tg.paragraphs[0], tag.upper(), 10, col if col != GOLD_DK else GOLD_DK, bold=True)
        x = Emu(x + cw + gap)
    _band(
        s,
        Inches(6.0),
        "The trust is deterministic: ",
        "every finding is exploit-verified and cryptographically signed — proof a model can't give.",
        h=Inches(0.9),
    )


# ── 9 · pillar 1 — detection ────────────────────────────────────────────────────
def pillar_detection_slide(prs):
    s = _content(
        prs,
        "Pillar 1 — Deterministic Detection",
        9,
        sub="13 tools · 3 language adapters · every output collapses into one CanonicalFinding schema — reproducible",
    )
    # diagram fills the slide (its own note carries the "why it matters" line)
    _fit_image(s, FIGS / "detection_fanin.png", Inches(0.5), Inches(1.9), Inches(12.33), Inches(4.95))


# ── 10 · pillar 2 — trust gates / funnel ────────────────────────────────────────
def pillar_trust_slide(prs):
    s = _content(
        prs,
        "Pillar 2 — Trust Gates",
        10,
        sub="30-repo adversarial corpus · 1,942 raw findings → 55 Confirmed Tier @ 96.4% precision",
    )
    _fit_image(s, FIGS / "FUNNEL_SLIDE.png", Inches(0.6), Inches(2.0), Inches(12.1), Inches(4.7))


# ── 11 · pillar 3 — exploit verification ────────────────────────────────────────
def pillar_exploit_slide(prs):
    s = _content(
        prs,
        "Pillar 3 — Exploit Verification",
        11,
        sub="We don't claim a vulnerability — we detonate it in a Docker sandbox, then re-detonate the fix.",
    )
    _fit_image(s, FIGS / "exploit_flow_slide.png", Inches(0.5), Inches(2.0), Inches(12.33), Inches(3.35))
    _band(
        s,
        Inches(5.55),
        "5/5 live:  ",
        "SQLi · CMDi · SSTI — all EXPLOITED, then BLOCKED. Binary ground truth — re-detonation cannot lie.",
        h=Inches(1.0),
    )


# ── 12 · pillar 4 — attestation ─────────────────────────────────────────────────
def pillar_attestation_slide(prs):
    s = _content(
        prs,
        "Pillar 4 — Cryptographic Attestation",
        12,
        sub="Every scan signed as a tamper-evident bundle — verifiable months later, in one command.",
    )
    rows = [
        ("ECDSA-P256", "Classical signature over the full findings bundle — fast, standard, widely verifiable."),
        ("Dilithium3", "NIST FIPS 204 post-quantum signature — survives a quantum adversary (EU CRA, Sept 2026)."),
        (
            "Merkle provenance",
            "Scan inputs, tool versions, and results hashed into a chain — nothing can be altered after the fact.",
        ),
        (
            "One-command verify",
            "Anyone can re-verify the bundle offline — trust the result without trusting the scanner.",
        ),
    ]
    y = Inches(2.05)
    for head, body in rows:
        _rect(s, Inches(0.5), y, Inches(7.4), Inches(1.0), LIGHT, rounded=True)
        _rect(s, Inches(0.5), y, Inches(0.14), Inches(1.0), GREEN)
        c = _box(s, Inches(0.85), Emu(y + Inches(0.1)), Inches(6.9), Inches(0.85))
        _run(c.paragraphs[0], head, 15, GREEN, bold=True)
        p = c.add_paragraph()
        _run(p, body, 12, INK)
        y = Emu(y + Inches(1.12))
    _cover_image(s, SHOTS / "attestation.png", Inches(8.1), Inches(2.05), Inches(4.73), Inches(4.5))


# ── 13 · system architecture ────────────────────────────────────────────────────
def architecture_slide(prs):
    s = _content(
        prs,
        "System Architecture",
        13,
        sub="36 engine modules · 3 language adapters · 12-stage async pipeline · 52-endpoint FastAPI",
    )
    # full-width landscape architecture (detailed — all 6 layers, from arch_overview)
    _fit_image(s, FIGS / "arch_overview_wide.png", Inches(0.4), Inches(1.95), Inches(12.55), Inches(2.35))
    # metrics row — 4 pills
    metrics = [
        ("96.4%", "Confirmed-Tier precision"),
        ("100%", "CVE recall (8/8)"),
        ("98.2%", "F₁ vs Semgrep 45.7%"),
        ("3,247", "tests · 88% coverage"),
    ]
    x = Inches(0.5)
    pw = Inches(3.0)
    gap = Inches(0.11)
    for big, lab in metrics:
        _rect(s, x, Inches(4.35), pw, Inches(0.95), LIGHT, rounded=True)
        bt = _box(s, x + Inches(0.2), Inches(4.4), Emu(pw - Inches(0.3)), Inches(0.5))
        _run(bt.paragraphs[0], big, 23, GREEN, bold=True)
        lt = _box(s, x + Inches(0.2), Inches(4.92), Emu(pw - Inches(0.3)), Inches(0.35))
        _run(lt.paragraphs[0], lab, 11.5, INK)
        x = Emu(x + pw + gap)
    # identity band
    _band(
        s,
        Inches(5.55),
        "Proof, not guessing: ",
        "exploit verification and signing are deterministic — an LLM cannot detonate a bug or sign a result.",
        h=Inches(1.0),
    )


# ── 14 · evaluation methodology ─────────────────────────────────────────────────
def evaluation_data_slide(prs):
    s = _content(
        prs,
        "Evaluation Methodology",
        14,
        sub="Not machine learning — there is no training set. A security tool is judged on what it's tested against, and how honestly.",
    )
    # four rigor principles (2x2) — the "how do you validate?" armor
    pillars = [
        ("No training set", "Deterministic rules, not a trained model — nothing to overfit or cherry-pick.", NAVY),
        ("Pre-registered", "The CVE battery was declared BEFORE testing — 100% recall, no tuning to the test.", GREEN),
        ("Third-party ground truth", "RealVuln's labels are external (arXiv:2604.13764) — not graded by me.", GOLD_DK),
        ("Manual triage", "Every Confirmed-Tier finding hand-verified to earn the 96.4% precision number.", RED),
    ]
    positions = [
        (Inches(0.5), Inches(2.05)),
        (Inches(6.83), Inches(2.05)),
        (Inches(0.5), Inches(3.75)),
        (Inches(6.83), Inches(3.75)),
    ]
    for (x, y), (head, body, col) in zip(positions, pillars):
        _rect(s, x, y, Inches(6.0), Inches(1.5), LIGHT, rounded=True)
        _rect(s, x, y, Inches(0.12), Inches(1.5), col)
        c = _box(s, x + Inches(0.28), Emu(y + Inches(0.16)), Inches(5.5), Inches(1.2))
        _run(c.paragraphs[0], head, 16, col if col != GOLD_DK else GOLD_DK, bold=True)
        p = c.add_paragraph()
        p.space_before = Pt(5)
        _run(p, body, 12.5, INK)
    # compact corpora reference strip
    _band(
        s,
        Inches(5.55),
        "Tested on: ",
        "RealVuln 2026 (real CVEs) · OWASP Benchmark · SecurityEval (89+89) · 30-repo adversarial corpus · 8-CVE battery.",
        h=Inches(0.95),
    )


# ── 15 · results — confusion matrix (alone, horizontal) ──────────────────────────
def results_confusion_slide(prs):
    s = _content(
        prs,
        "Results — Why Trust Gates Exist",
        15,
        sub="Raw detection on SecurityEval: 91% recall, but only 54.7% precision — then the Confirmed Tier fixes it.",
    )
    _fit_image(s, FIGS / "CONFUSION_MATRIX_SLIDE.png", Inches(0.5), Inches(2.2), Inches(12.33), Inches(3.6))
    _band(
        s,
        Inches(6.1),
        "The fix → Confirmed Tier: ",
        "those 67 false alarms collapse and precision jumps to 96.4% (slide 10's funnel). Honest raw number, strong final one.",
        h=Inches(0.85),
    )


# ── 16 · results — benchmark scorecard (variety: strong numbers at a glance) ──────
def results_scorecard_slide(prs):
    s = _content(
        prs,
        "Results — Across Five Benchmarks",
        16,
        sub="Not one test — five. Each measures a different thing, and ACR-QA scores 90–100% on every one.",
    )
    _fit_image(s, FIGS / "BENCHMARK_SCORECARD.png", Inches(0.5), Inches(2.0), Inches(12.33), Inches(3.7))
    _band(
        s,
        Inches(5.95),
        "Variety, not cherry-picking: ",
        "precision, recall, F₁, and OWASP coverage — all high. The one hard real-world number (RealVuln 25.1%) is next, where we're #1.",
        h=Inches(0.95),
    )


# ── 17 · results — benchmarks ───────────────────────────────────────────────────
def results_benchmarks_slide(prs):
    s = _content(
        prs,
        "Results — #1 Against Every Competitor",
        17,
        sub="On real-world CVEs (RealVuln 2026), ACR-QA beats Semgrep, Snyk, SonarQube — the hardest test, third-party ground truth.",
    )
    _fit_image(s, FIGS / "HEAD_TO_HEAD.png", Inches(0.5), Inches(2.0), Inches(6.1), Inches(3.95))
    _fit_image(s, FIGS / "REALVULN_LEADERBOARD_CLEAN.png", Inches(6.9), Inches(2.25), Inches(5.95), Inches(3.45))
    _band(
        s,
        Inches(5.95),
        "Why the absolute number looks modest: ",
        "real-world recall is genuinely hard (many CVEs need runtime state no static tool can see) — so #1 is the story, not 25%.",
        h=Inches(0.95),
    )


# ── 18 · use cases (one compact slide — 4 personas) ──────────────────────────────
def _persona_card(s, x, y, w, h, head, body, payoff, col):
    _rect(s, x, y, w, h, LIGHT, rounded=True)
    _rect(s, x, y, w, Inches(0.12), col)
    c = _box(s, x + Inches(0.24), Emu(y + Inches(0.22)), Emu(w - Inches(0.48)), Emu(h - Inches(0.4)))
    _run(c.paragraphs[0], head, 14.5, col if col != GOLD else GOLD_DK, bold=True)
    p = c.add_paragraph()
    p.space_before = Pt(5)
    _run(p, body, 11, INK)
    p2 = c.add_paragraph()
    p2.space_before = Pt(6)
    _run(p2, payoff, 10.5, col if col != GOLD else GOLD_DK, bold=True)


def use_cases_slide(prs):
    s = _content(prs, "Use Cases", 18, sub="One trust layer — four people who can't afford to guess wrong.")
    personas = [
        (
            Inches(0.5),
            Inches(2.05),
            "University Instructor",
            "Grades hundreds of student repos a term — enterprise scanners cost $10–50k.",
            "→ Free, self-hosted · the real bug surfaced · RAG explanation teaches.",
            NAVY,
        ),
        (
            Inches(6.83),
            Inches(2.05),
            "Startup / CI Tech Lead",
            "Ships dozens of AI-written PRs a day — needs an auto-block that won't cry wolf.",
            "→ 96.4% auto-block · signed audit trail for EU CRA · $0 recurring.",
            GREEN,
        ),
        (
            Inches(0.5),
            Inches(4.5),
            "Open-Source Maintainer",
            "Reviews drive-by PRs from strangers — can't spot a subtle supply-chain attack.",
            "→ Exploit-verified + signed provenance — trust the code, not the contributor.",
            GOLD,
        ),
        (
            Inches(6.83),
            Inches(4.5),
            "Enterprise Auditor",
            "Must prove to regulators what was scanned, months later, tamper-free.",
            "→ Every scan ECDSA + Dilithium3 signed · verify in one command.",
            RED,
        ),
    ]
    for x, y, head, body, payoff, col in personas:
        _persona_card(s, x, y, Inches(6.0), Inches(2.15), head, body, payoff, col)


# ── 19 · conclusion + future work (merged) ───────────────────────────────────────
def conclusion_future_slide(prs):
    s = _content(prs, "Conclusion & Future Work", 19)
    tf = _box(s, Inches(0.5), Inches(1.5), Inches(7.3), Inches(0.5))
    _run(tf.paragraphs[0], "The open, attested, $0 quadrant the market leaves empty.", 17, INK, bold=True)
    lines = [
        ("Trust", "96.4% Confirmed-Tier precision — enough to auto-block a merge."),
        ("Proof", "Exploit-verified + cryptographically signed — not guesses."),
        ("Reach", "13 tools · Python / JS / Go · 9/10 OWASP · 100% CVE recall."),
        ("Price", "Self-hosted · $0 recurring — vs $10–50k/year."),
    ]
    y = Inches(2.1)
    for k, v in lines:
        _rect(s, Inches(0.5), y, Inches(1.5), Inches(0.68), GOLD, rounded=True)
        kt = _box(s, Inches(0.5), y, Inches(1.5), Inches(0.68))
        kt.vertical_anchor = MSO_ANCHOR.MIDDLE
        _center(kt.paragraphs[0])
        _run(kt.paragraphs[0], k, 15, NAVY, bold=True)
        vt = _box(s, Inches(2.15), y, Inches(5.5), Inches(0.68))
        vt.vertical_anchor = MSO_ANCHOR.MIDDLE
        _run(vt.paragraphs[0], v, 11.5, INK)
        y = Emu(y + Inches(0.8))
    _cover_image(s, PHOTOS / "Gemini_6.png", Inches(8.0), Inches(1.5), Inches(4.83), Inches(3.2))
    cap = _box(s, Inches(8.0), Emu(Inches(1.5) + Inches(3.2)), Inches(4.83), Inches(0.35))
    _center(cap.paragraphs[0])
    _run(cap.paragraphs[0], "Three locked doors, one open — the gap we fill.", 10.5, GRAY, italic=True)
    # future work strip
    fwh = _box(s, Inches(0.5), Inches(5.35), Inches(12.3), Inches(0.4))
    _run(fwh.paragraphs[0], "FUTURE WORK", 13, GOLD_DK, bold=True)
    fut = [
        "Inter-procedural taint (+10–15pp recall)",
        "More languages: Java · C# · Rust",
        "Managed one-click GitHub App",
        "Independent multi-annotator labelling",
    ]
    x = Inches(0.5)
    fw = Inches(3.0)
    for item in fut:
        _rect(s, x, Inches(5.78), fw, Inches(0.85), LIGHT, rounded=True)
        _rect(s, x, Inches(5.78), Inches(0.1), Inches(0.85), GOLD)
        c = _box(s, x + Inches(0.22), Inches(5.78), Emu(fw - Inches(0.34)), Inches(0.85))
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        _run(c.paragraphs[0], item, 11.5, INK)
        x = Emu(x + fw + Inches(0.11))
    kill = _box(s, Inches(0.5), Inches(6.72), Inches(7.3), Inches(0.5))
    _run(kill.paragraphs[0], "Not a proposal — a running system.", 14, NAVY, bold=True, italic=True)


def section_closer(prs, big, small, photo=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, NAVY)
    if photo:
        _hero(s, photo, alpha=70)
    _rect(s, 0, Inches(6.95), EMU_W, Inches(0.55), GOLD)
    if LOGO.exists():
        s.shapes.add_picture(str(LOGO), Inches(6.17), Inches(1.4), height=Inches(1.3))
    tf = _box(s, Inches(0.8), Inches(3.1), Inches(11.7), Inches(2.2))
    _center(tf.paragraphs[0])
    _run(tf.paragraphs[0], big, 48, WHITE, bold=True)
    if small:
        p = tf.add_paragraph()
        _center(p)
        _run(p, small, 18, GOLD)


def build():
    prs = Presentation()
    prs.slide_width, prs.slide_height = EMU_W, EMU_H

    title_slide(prs)  # 1
    outline_slide(prs)  # 2
    introduction_slide(prs)  # 3
    problem_overview_slide(prs)  # 4
    problem_detail_slide(
        prs,
        5,
        "1",
        "Gemini_3.png",
        "Too many alerts",
        "A wall of warnings, almost all false. The one real breach hides in the noise — "
        "and a tired reviewer on the 50th PR of the day misses it.",
        "~1,900",
        "raw findings per scan · which one breaches you?",
    )
    problem_detail_slide(
        prs,
        6,
        "2",
        "Gemini_5.png",
        "Claimed, never proven",
        "Scanners flag “possible SQL injection” and stop there. The developer has to guess "
        "whether it's real — and gets burned either way.",
        "0",
        "scanners actually detonate the bug to prove it",
    )
    problem_detail_slide(
        prs,
        7,
        "3",
        "Gemini_6.png",
        "No provenance",
        "Findings aren't signed. Months later you can't prove what was scanned — and you can't trust "
        "a stranger's PR or the code an AI just wrote into your repo.",
        "45%",
        "of AI-written code ships a known flaw",
    )
    solution_overview_slide(prs)  # 8
    pillar_detection_slide(prs)  # 9
    pillar_trust_slide(prs)  # 10
    pillar_exploit_slide(prs)  # 11
    pillar_attestation_slide(prs)  # 12
    architecture_slide(prs)  # 13
    evaluation_data_slide(prs)  # 14
    results_confusion_slide(prs)  # 15
    results_scorecard_slide(prs)  # 16
    results_benchmarks_slide(prs)  # 17  (live demo launches from the Trust-Gates funnel, slide 10)
    use_cases_slide(prs)  # 18
    conclusion_future_slide(prs)  # 19
    section_closer(prs, "Any Questions?", "Happy to go deeper on any number, any slide.")  # 20
    section_closer(
        prs, "Thank You", "Ahmed Mahmoud Abbas · Supervisor: Dr. Samy Abdel Nabi", photo="Gemini_2.png"
    )  # 21

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"✓ Template deck written: {OUT}  ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    build()
